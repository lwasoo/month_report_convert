"""PowerPoint template inspection and slide-level operations."""

from __future__ import annotations

from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt

from ..common import log, normalize_text
from ..models import TemplateSlide


def extract_template_slides(prs: Presentation) -> list[TemplateSlide]:
    log("读取 PPT 模板结构")
    slides: list[TemplateSlide] = []
    for idx, slide in enumerate(prs.slides, start=1):
        if idx == 1 or idx == len(prs.slides):
            continue
        title = ""
        for shp in slide.shapes:
            if getattr(shp, "has_text_frame", False) and normalize_text(getattr(shp, "text", "")):
                if shp.top <= Inches(1.6):
                    title = normalize_text(shp.text)
                    break
        has_table = any(getattr(shp, "has_table", False) for shp in slide.shapes)
        slides.append(TemplateSlide(idx, title or f"第{idx}页", has_table))
    log(f"PPT 模板解析完成: 需要填充 {len(slides)} 页")
    return slides


def remove_auto_shapes(slide) -> None:
    for shp in list(slide.shapes):
        if shp.name.startswith("AUTO_CONTENT_") or shp.name.startswith("AUTO_LAYOUT_"):
            slide.shapes._spTree.remove(shp._element)


def insert_slide_after(prs: Presentation, after_index: int, layout) -> Any:
    new_slide = prs.slides.add_slide(layout)
    sld_id_lst = prs.slides._sldIdLst  # type: ignore[attr-defined]
    new_id = sld_id_lst[-1]
    sld_id_lst.remove(new_id)
    sld_id_lst.insert(after_index + 1, new_id)
    return prs.slides[after_index + 1]


def set_slide_title_text(slide, title: str) -> None:
    target = None
    for shp in slide.shapes:
        if not getattr(shp, "has_text_frame", False):
            continue
        txt = normalize_text(getattr(shp, "text", ""))
        if shp.top <= Inches(1.6):
            target = shp
            if txt:
                break
    if target is None:
        box = slide.shapes.add_textbox(Inches(0.55), Inches(0.25), Inches(8.5), Inches(0.6))
        box.name = "AUTO_OVERFLOW_TITLE"
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.name = "Microsoft YaHei"
        return
    target.text = title
