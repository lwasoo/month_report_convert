"""Formal visual layout variants for generated slide content."""

from __future__ import annotations

import random

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches

from .content_regions import (
    _find_content_regions,
    _write_lines_to_box,
    _write_rows_across_regions,
    get_theme_palette,
    split_balanced,
    suggest_font_size,
)
from ..models import SlideDraft


def choose_formal_variant(title: str, bullets: list[str], diversity: str, seed: int, slide_index: int) -> str:
    title = title or ""
    if diversity == "none":
        return "single_column"
    if len(bullets) >= 6:
        return "two_column"
    if any(k in title for k in ["风险", "诉讼", "仲裁"]):
        base = ["risk_matrix", "timeline", "two_column"]
    elif any(k in title for k in ["数据", "统计", "流程"]):
        base = ["kpi_cards", "two_column", "single_column"]
    elif any(k in title for k in ["项目", "进度", "政策", "合规"]):
        base = ["timeline", "two_column", "single_column"]
    else:
        base = ["two_column", "single_column", "kpi_cards"]

    if len(bullets) <= 2:
        return "single_column"
    if len(bullets) >= 4 and diversity == "high":
        base = ["two_column", "timeline", "kpi_cards", "risk_matrix"]
    elif diversity == "low":
        base = [base[0]]

    rng = random.Random((seed or 0) * 131 + slide_index * 17 + len(bullets))
    return base[rng.randrange(len(base))]


def add_visual_accent(slide, slide_index: int, palette: dict[str, RGBColor]) -> None:
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.52), Inches(1.16), Inches(0.06), Inches(3.9))
    accent.name = f"AUTO_LAYOUT_ACCENT_{slide_index}"
    accent.fill.solid()
    accent.fill.fore_color.rgb = palette["accent"]
    accent.line.fill.background()


def render_single_column(slide, slide_index: int, rows: list[str], font_size: int) -> None:
    _write_lines_to_box(slide, f"AUTO_CONTENT_{slide_index}", Inches(0.72), Inches(1.12), Inches(8.55), Inches(5.75), rows, font_size, bullet="• ")


def render_two_column(slide, slide_index: int, rows: list[str], font_size: int) -> None:
    left_rows, right_rows = split_balanced(rows)
    col_top = Inches(1.12)
    col_height = Inches(5.75)
    col_width = Inches(4.18)
    _write_lines_to_box(slide, f"AUTO_CONTENT_{slide_index}_L", Inches(0.72), col_top, col_width, col_height, left_rows, font_size, bullet="• ")
    _write_lines_to_box(slide, f"AUTO_CONTENT_{slide_index}_R", Inches(5.08), col_top, col_width, col_height, right_rows, font_size, bullet="• ")


def render_timeline(slide, slide_index: int, rows: list[str], font_size: int, palette: dict[str, RGBColor]) -> None:
    start_top = Inches(1.18)
    step = Inches(0.96)
    for idx, row in enumerate(rows[:6]):
        y = start_top + idx * step
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(0.92), y + Inches(0.05), Inches(0.12), Inches(0.12))
        dot.name = f"AUTO_LAYOUT_DOT_{slide_index}_{idx}"
        dot.fill.solid()
        dot.fill.fore_color.rgb = palette["accent"]
        dot.line.fill.background()
        if idx < len(rows) - 1:
            line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.97), y + Inches(0.16), Inches(0.02), Inches(0.6))
            line.name = f"AUTO_LAYOUT_LINE_{slide_index}_{idx}"
            line.fill.solid()
            line.fill.fore_color.rgb = palette["line"]
            line.line.fill.background()
        _write_lines_to_box(slide, f"AUTO_CONTENT_{slide_index}_{idx}", Inches(1.2), y, Inches(7.9), Inches(0.68), [row], font_size, bullet="")


def render_kpi_cards(slide, slide_index: int, rows: list[str], font_size: int, palette: dict[str, RGBColor]) -> None:
    rows = rows[:4]
    card_w = Inches(4.1)
    card_h = Inches(1.0)
    positions = [(Inches(0.72), Inches(1.22)), (Inches(4.95), Inches(1.22)), (Inches(0.72), Inches(2.45)), (Inches(4.95), Inches(2.45))]
    for idx, row in enumerate(rows):
        left, top = positions[idx]
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, card_w, card_h)
        card.name = f"AUTO_LAYOUT_CARD_{slide_index}_{idx}"
        card.fill.solid()
        card.fill.fore_color.rgb = palette["soft"]
        card.line.color.rgb = palette["line"]
        _write_lines_to_box(slide, f"AUTO_CONTENT_{slide_index}_{idx}", left + Inches(0.18), top + Inches(0.12), card_w - Inches(0.32), card_h - Inches(0.24), [row], font_size, bullet="")


def render_risk_matrix(slide, slide_index: int, rows: list[str], font_size: int, palette: dict[str, RGBColor]) -> None:
    left = Inches(0.86)
    top = Inches(1.24)
    width = Inches(8.1)
    height = Inches(5.2)
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
    box.name = f"AUTO_LAYOUT_RISKBOX_{slide_index}"
    box.fill.solid()
    box.fill.fore_color.rgb = palette["soft"]
    box.line.color.rgb = palette["line"]
    for idx, row in enumerate(rows[:6]):
        _write_lines_to_box(slide, f"AUTO_CONTENT_{slide_index}_{idx}", left + Inches(0.22), top + Inches(0.18 + idx * 0.82), width - Inches(0.44), Inches(0.62), [row], font_size, bullet="• ")


def add_formal_layout_content_v2(
    slide,
    draft: SlideDraft,
    has_table: bool,
    title: str,
    theme: str,
    diversity: str,
    seed: int,
    slide_width: int | None = None,
    slide_height: int | None = None,
) -> bool:
    if not draft.bullets:
        return False

    if has_table:
        if slide_width is None or slide_height is None:
            slide_width = int(Inches(10))
            slide_height = int(Inches(7.5))
        rows = draft.bullets
        font_size = suggest_font_size(rows, base=13)
        regions = _find_content_regions(slide, slide_width, slide_height)
        _write_rows_across_regions(slide, f"AUTO_CONTENT_{draft.slide_index}", rows, regions, font_size, "• ")
        return True

    palette = get_theme_palette(theme)
    rows = draft.bullets

    add_visual_accent(slide, draft.slide_index, palette)
    font_size = suggest_font_size(rows, base=15)
    variant = choose_formal_variant(title, rows, diversity, seed, draft.slide_index)

    if variant == "two_column":
        render_two_column(slide, draft.slide_index, rows, font_size)
    elif variant == "timeline":
        render_timeline(slide, draft.slide_index, rows, font_size, palette)
    elif variant == "kpi_cards":
        render_kpi_cards(slide, draft.slide_index, rows, font_size, palette)
    elif variant == "risk_matrix":
        render_risk_matrix(slide, draft.slide_index, rows, font_size, palette)
    else:
        render_single_column(slide, draft.slide_index, rows, font_size)
    return False


