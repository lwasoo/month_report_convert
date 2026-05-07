"""Content box placement inside available slide regions."""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from ..common import normalize_text
from ..models import SlideDraft


def style_paragraph(paragraph, size: int) -> None:
    paragraph.font.name = "Microsoft YaHei"
    paragraph.font.size = Pt(size)
    paragraph.alignment = PP_ALIGN.LEFT
    paragraph.space_after = Pt(2)
    paragraph.space_before = Pt(0)
    paragraph.line_spacing = 1.15


def get_theme_palette(theme: str) -> dict[str, RGBColor]:
    if theme == "corporate_gray":
        return {
            "accent": RGBColor(83, 108, 138),
            "soft": RGBColor(230, 235, 240),
            "line": RGBColor(190, 198, 208),
            "text": RGBColor(48, 72, 96),
        }
    if theme == "legal_red":
        return {
            "accent": RGBColor(161, 47, 47),
            "soft": RGBColor(245, 233, 233),
            "line": RGBColor(219, 188, 188),
            "text": RGBColor(88, 46, 46),
        }
    return {
        "accent": RGBColor(28, 84, 140),
        "soft": RGBColor(232, 240, 248),
        "line": RGBColor(184, 207, 228),
        "text": RGBColor(34, 63, 92),
    }


def suggest_font_size(lines: list[str], base: int = 16) -> int:
    total = sum(len(x) for x in lines)
    longest = max((len(x) for x in lines), default=0)
    size = base
    if total > 260 or longest > 64:
        size -= 1
    if total > 360 or longest > 82:
        size -= 1
    if total > 460 or longest > 100:
        size -= 1
    return max(size, 12)


def split_balanced(lines: list[str]) -> tuple[list[str], list[str]]:
    if len(lines) <= 1:
        return lines, []
    mid = (len(lines) + 1) // 2
    return lines[:mid], lines[mid:]


def _write_lines_to_box(slide, name: str, left, top, width, height, rows: list[str], font_size: int, bullet: str) -> None:
    if not rows:
        return
    box = slide.shapes.add_textbox(left, top, width, height)
    box.name = name
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, row in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet}{row}" if bullet else row
        style_paragraph(p, size=font_size)
        p.level = 0


def _subtract_interval(intervals: list[tuple[int, int]], block_start: int, block_end: int) -> list[tuple[int, int]]:
    out: list[tuple[int, int]] = []
    for start, end in intervals:
        if block_end <= start or block_start >= end:
            out.append((start, end))
            continue
        if block_start > start:
            out.append((start, block_start))
        if block_end < end:
            out.append((block_end, end))
    return out


def _detect_title_bottom(slide, default_top: int) -> int:
    bottom = default_top
    for shp in slide.shapes:
        if not getattr(shp, "has_text_frame", False):
            continue
        txt = normalize_text(getattr(shp, "text", ""))
        if not txt:
            continue
        if shp.top <= Inches(1.6):
            candidate = int(shp.top + shp.height + Inches(0.05))
            if candidate > bottom:
                bottom = candidate
    return bottom


def _find_content_regions(slide, slide_width: int, slide_height: int) -> list[tuple[int, int, int, int]]:
    left = int(Inches(0.55))
    right = int(slide_width - Inches(0.55))
    top_limit = _detect_title_bottom(slide, int(Inches(1.05)))
    bottom_limit = int(slide_height - Inches(0.25))
    intervals: list[tuple[int, int]] = [(top_limit, bottom_limit)]

    for shp in slide.shapes:
        if getattr(shp, "has_table", False):
            t0 = int(max(shp.top - Inches(0.05), top_limit))
            t1 = int(min(shp.top + shp.height + Inches(0.05), bottom_limit))
            intervals = _subtract_interval(intervals, t0, t1)

    intervals = [(start, end) for start, end in intervals if end - start >= int(Inches(0.45))]
    if not intervals:
        return [(left, top_limit, right - left, int(max(bottom_limit - top_limit, Inches(0.6))))]
    return [(left, start, right - left, end - start) for start, end in intervals]


def _find_content_region(slide, slide_width: int, slide_height: int) -> tuple[int, int, int, int]:
    regions = _find_content_regions(slide, slide_width, slide_height)
    return max(regions, key=lambda item: item[3])


def _write_rows_across_regions(
    slide,
    base_name: str,
    rows: list[str],
    regions: list[tuple[int, int, int, int]],
    font_size: int,
    bullet: str,
) -> None:
    if not rows or not regions:
        return

    total_height = sum(region[3] for region in regions)
    start = 0
    remaining_rows = len(rows)
    remaining_regions = len(regions)

    for idx, (left, top, width, height) in enumerate(regions):
        if remaining_rows <= 0:
            break
        if idx == len(regions) - 1:
            chunk = rows[start:]
        else:
            ratio = height / max(total_height, 1)
            take = max(1, round(len(rows) * ratio))
            take = min(take, remaining_rows - (remaining_regions - 1))
            chunk = rows[start : start + take]
        _write_lines_to_box(slide, f"{base_name}_{idx}", left, top, width, height, chunk, font_size, bullet)
        start += len(chunk)
        remaining_rows -= len(chunk)
        remaining_regions -= 1


def add_content_textbox(slide, draft: SlideDraft, has_table: bool, slide_width: int, slide_height: int) -> None:
    if not draft.bullets:
        return

    rows = draft.bullets
    font_size = suggest_font_size(rows, base=15 if has_table else 16)
    if has_table:
        regions = _find_content_regions(slide, slide_width, slide_height)
        _write_rows_across_regions(slide, f"AUTO_CONTENT_{draft.slide_index}", rows, regions, font_size, "• ")
        return

    left, top, width, height = _find_content_region(slide, slide_width, slide_height)
    if height < int(Inches(0.45)):
        return
    _write_lines_to_box(slide, f"AUTO_CONTENT_{draft.slide_index}", left, top, width, height, rows, font_size, bullet="• ")


