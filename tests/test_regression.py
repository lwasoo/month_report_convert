from __future__ import annotations

import json
import sys
import unittest
from pathlib import Path
from tempfile import TemporaryDirectory
from unittest.mock import patch

from docx import Document
from pptx import Presentation

from doc_sanitizer.document_io import apply_mapping_to_file, collect_texts_for_path, restore_file
from doc_sanitizer.fuzzy_mapping import (
    build_external_ai_prompt_sections,
    closest_placeholder_for_token,
    find_placeholder_like_tokens,
    payload_from_json_text,
    placeholder_token_category,
    suggest_placeholder_repairs,
    unresolved_placeholder_tokens,
)
from doc_sanitizer.llm_assist import build_llm_candidate_contexts, chunk_texts_for_llm, count_texts_with_existing_terms, select_texts_for_llm, stable_text_hash
from doc_sanitizer.mapping import mapping_entries, read_mapping
from report_converter.common import log, route_logs_to
from gui_app.self_update import build_windows_update_script, windows_update_target_path
from gui_app.update_checker import (
    ReleaseAsset,
    UpdateInfo,
    compare_versions,
    download_release_asset,
    sanitize_filename,
    unique_path,
)


class MappingAndPromptTests(unittest.TestCase):
    def test_prompt_copy_section_excludes_sensitive_originals(self) -> None:
        # 通过用例：外部 AI 可复制区只能包含占位符，不能泄露原始敏感词。
        # 内部审核区可以保留原始名称线索，供用户自己判断归组是否合理。
        payload = payload_from_json_text(
            json.dumps(
                {
                    "entries": [
                        {
                            "placeholder": "__COMPANY_001__",
                            "original": "Acme",
                            "category": "COMPANY",
                            "enabled": True,
                        },
                        {
                            "placeholder": "__COMPANY_002__",
                            "original": "Acme Technology Co Ltd",
                            "category": "COMPANY",
                            "enabled": True,
                        },
                    ]
                },
                ensure_ascii=False,
            )
        )

        prompt, audit = build_external_ai_prompt_sections(payload)

        self.assertIn("__COMPANY_002__ / __COMPANY_001__", prompt)
        self.assertIn("不得新增", prompt)
        self.assertIn("PROJECT_015", prompt)
        self.assertNotIn("Acme", prompt)
        self.assertNotIn("映射摘要", prompt)
        self.assertNotIn("人工确认", prompt)
        self.assertIn("Acme Technology Co Ltd", audit)

    def test_placeholder_repair_scores_split_auto_and_confirmation_cases(self) -> None:
        # 通过用例：轻微损坏的占位符应进入高置信自动修复；
        # 缺字/编号简写这类更可疑的写法只进入“需要用户确认”的分数段。
        payload = payload_from_json_text(
            json.dumps(
                {
                    "entries": [
                        {
                            "placeholder": "__COMPANY_001__",
                            "original": "Acme",
                            "category": "COMPANY",
                            "enabled": True,
                        }
                    ]
                },
                ensure_ascii=False,
            )
        )

        repairs = suggest_placeholder_repairs("Keep COMPANY_001 and maybe COMY_01.", mapping_entries(payload), min_score=0.70)
        scores = {repair.token: repair.score for repair in repairs}

        self.assertGreaterEqual(scores["COMPANY_001"], 0.90)
        self.assertGreaterEqual(scores["COMY_01"], 0.70)
        self.assertLess(scores["COMY_01"], 0.90)

    def test_placeholder_repair_prefers_matching_index_over_category_similarity(self) -> None:
        # 失败用例：历史问题是前缀很像时会把编号不同的占位符误配过去。
        # 这里要求编号优先：MPANY_007 应配 __COMPANY_007__，不能配 __COMPANY_001__。
        payload = payload_from_json_text(
            json.dumps(
                {
                    "entries": [
                        {
                            "placeholder": "__COMPANY_001__",
                            "original": "Acme",
                            "category": "COMPANY",
                            "enabled": True,
                        },
                        {
                            "placeholder": "__COMPANY_007__",
                            "original": "Example Holdings",
                            "category": "COMPANY",
                            "enabled": True,
                        },
                    ]
                },
                ensure_ascii=False,
            )
        )

        repairs = suggest_placeholder_repairs("MPANY_007", mapping_entries(payload), min_score=0.70)
        self.assertEqual(repairs[0].canonical, "__COMPANY_007__")
        self.assertGreaterEqual(repairs[0].score, 0.90)

    def test_placeholder_repair_rejects_different_index_with_same_category(self) -> None:
        # 失败用例：PROJECT_015 不能因为类别相同就被猜成 __PROJECT_005__。
        payload = payload_from_json_text(
            json.dumps(
                {
                    "entries": [
                        {
                            "placeholder": "__PROJECT_005__",
                            "original": "Project A",
                            "category": "PROJECT",
                            "enabled": True,
                        }
                    ]
                },
                ensure_ascii=False,
            )
        )

        repairs = suggest_placeholder_repairs("PROJECT_015", mapping_entries(payload), min_score=0.70)
        self.assertEqual(repairs, [])

    def test_placeholder_repair_confirms_extra_digit_but_does_not_auto_accept(self) -> None:
        # 失败/边界用例：_CODE_0077_ 和 __CODE_007__ 很像，但编号多了一位，
        # 应进入确认分数段，而不是自动修复。
        payload = payload_from_json_text(
            json.dumps(
                {
                    "entries": [
                        {
                            "placeholder": "__CODE_007__",
                            "original": "BU11",
                            "category": "CODE",
                            "enabled": True,
                        }
                    ]
                },
                ensure_ascii=False,
            )
        )

        repairs = suggest_placeholder_repairs("_CODE_0077_", mapping_entries(payload), min_score=0.70)
        self.assertEqual(repairs[0].canonical, "__CODE_007__")
        self.assertGreaterEqual(repairs[0].score, 0.70)
        self.assertLess(repairs[0].score, 0.90)

    def test_unresolved_placeholder_tokens_report_unknown_ai_ids(self) -> None:
        # 失败用例：外部 AI 可能生成映射表不存在的新编号；不能静默显示还原完成。
        payload = payload_from_json_text(
            json.dumps(
                {
                    "entries": [
                        {
                            "placeholder": "__CODE_007__",
                            "original": "BU11-1件",
                            "category": "CODE",
                            "enabled": True,
                        }
                    ]
                },
                ensure_ascii=False,
            )
        )

        tokens = unresolved_placeholder_tokens(
            "已还原 BU11-1件，但 _CODE_008_ 和 __SUPPLIER_004T__ 仍残留。",
            mapping_entries(payload),
        )

        self.assertEqual(tokens, ["_CODE_008_", "__SUPPLIER_004T__"])
        self.assertEqual(placeholder_token_category("__SUPPLIER_004T__"), "SUPPLIER")

    def test_unknown_placeholder_prefills_closest_numbered_candidate(self) -> None:
        # 通过用例：未知占位符也应先按类别和编号靠近程度给出预填建议。
        payload = payload_from_json_text(
            json.dumps(
                {
                    "entries": [
                        {
                            "placeholder": "__SUPPLIER_001__",
                            "original": "AmpLink",
                            "category": "SUPPLIER",
                            "enabled": True,
                        },
                        {
                            "placeholder": "__SUPPLIER_003__",
                            "original": "苏州立讯技术",
                            "category": "SUPPLIER",
                            "enabled": True,
                        },
                    ]
                },
                ensure_ascii=False,
            )
        )

        placeholder, score = closest_placeholder_for_token("__SUPPLIER_003T__", mapping_entries(payload), min_score=0.65)

        self.assertEqual(placeholder, "__SUPPLIER_003__")
        self.assertGreaterEqual(score, 0.90)

    def test_broken_placeholder_tail_is_detected_inside_joined_text(self) -> None:
        # 失败用例：外部 AI 可能把 __COMPANY_007__ 改成 QuantaMPANY_007 这种粘连残片。
        tokens = find_placeholder_like_tokens("客户CM厂QuantaMPANY_007")

        self.assertIn("MPANY_007", tokens)

    def test_repair_does_not_partially_replace_embedded_placeholder(self) -> None:
        # 失败用例：不能先把 __SUPPLIER_003T__ 的子串 __SUPPLIER_003 替换掉，
        # 否则会产生 __SUPPLIER_003__T__ 这种更坏的残留。
        payload = payload_from_json_text(
            json.dumps(
                {
                    "entries": [
                        {
                            "placeholder": "__SUPPLIER_003__",
                            "original": "苏州立讯技术",
                            "category": "SUPPLIER",
                            "enabled": True,
                        }
                    ]
                },
                ensure_ascii=False,
            )
        )

        repairs = suggest_placeholder_repairs("供应商 __SUPPLIER_003T__ 已签约", mapping_entries(payload), min_score=0.70)

        self.assertEqual(repairs, [])

    def test_repair_still_handles_placeholder_next_to_chinese_text(self) -> None:
        # 通过用例：中文前后缀不算英文粘连，仍应识别 _CODE_007_ 这类轻微损坏占位符。
        payload = payload_from_json_text(
            json.dumps(
                {
                    "entries": [
                        {
                            "placeholder": "__CODE_007__",
                            "original": "BU11-1件",
                            "category": "CODE",
                            "enabled": True,
                        }
                    ]
                },
                ensure_ascii=False,
            )
        )

        repairs = suggest_placeholder_repairs("其中_CODE_007_类", mapping_entries(payload), min_score=0.70)

        self.assertEqual(repairs[0].canonical, "__CODE_007__")

    def test_invalid_mapping_json_is_rejected(self) -> None:
        # 失败用例：外部输入不是映射对象时应直接报错，避免生成误导性的 Prompt。
        with self.assertRaises(ValueError):
            payload_from_json_text(json.dumps(["not", "a", "mapping"]))

        # 失败用例：JSON 结构存在但没有任何有效 entries/replacements，也应报错。
        with self.assertRaises(ValueError):
            payload_from_json_text(json.dumps({"entries": []}))

    def test_very_weak_placeholder_match_is_not_suggested(self) -> None:
        # 失败用例：太不像的 token 不应被猜测成某个占位符，避免错误还原。
        payload = payload_from_json_text(
            json.dumps(
                {
                    "entries": [
                        {
                            "placeholder": "__COMPANY_001__",
                            "original": "Acme",
                            "category": "COMPANY",
                            "enabled": True,
                        }
                    ]
                },
                ensure_ascii=False,
            )
        )

        repairs = suggest_placeholder_repairs("This token ABC_99 should not match.", mapping_entries(payload), min_score=0.70)
        self.assertEqual(repairs, [])


class DocumentRoundTripTests(unittest.TestCase):
    def test_docx_sanitize_and_restore_with_confirmed_placeholder_repairs(self) -> None:
        # 通过用例：覆盖最核心的 docx 链路：
        # 原文 -> 脱敏文档/映射 JSON -> 外部 AI 改坏占位符 -> 按确认结果还原。
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source_path = root / "source.docx"
            sanitized_path = root / "sanitized.docx"
            edited_path = root / "edited.docx"
            restored_path = root / "restored.docx"
            mapping_path = root / "mapping.json"

            doc = Document()
            doc.add_paragraph("Acme Technology Co Ltd signed Project Phoenix.")
            doc.save(source_path)

            payload = {
                "version": 2,
                "entries": [
                    {
                        "placeholder": "__COMPANY_001__",
                        "original": "Acme Technology Co Ltd",
                        "category": "COMPANY",
                        "enabled": True,
                        "source": "test",
                    },
                    {
                        "placeholder": "__PROJECT_001__",
                        "original": "Project Phoenix",
                        "category": "PROJECT",
                        "enabled": True,
                        "source": "test",
                    },
                ],
            }

            apply_mapping_to_file(source_path, sanitized_path, payload, mapping_path)
            sanitized_text = "\n".join(p.text for p in Document(sanitized_path).paragraphs)
            self.assertIn("__COMPANY_001__", sanitized_text)
            self.assertIn("__PROJECT_001__", sanitized_text)
            self.assertEqual(read_mapping(mapping_path)["replacements"]["__COMPANY_001__"], "Acme Technology Co Ltd")

            edited = Document()
            edited.add_paragraph("COMPANY_001 signed __PROJECT-001__.")
            edited.save(edited_path)

            restore_file(
                input_path=edited_path,
                output_path=restored_path,
                mapping_path=mapping_path,
                placeholder_repairs={
                    "COMPANY_001": "__COMPANY_001__",
                    "__PROJECT-001__": "__PROJECT_001__",
                },
            )
            restored_text = "\n".join(p.text for p in Document(restored_path).paragraphs)
            self.assertIn("Acme Technology Co Ltd", restored_text)
            self.assertIn("Project Phoenix", restored_text)

    def test_restore_rejects_empty_mapping_file(self) -> None:
        # 失败用例：映射 JSON 没有有效条目时，不能生成一个看似成功但未还原的文件。
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            input_path = root / "input.docx"
            output_path = root / "output.docx"
            mapping_path = root / "mapping.json"

            doc = Document()
            doc.add_paragraph("__COMPANY_001__")
            doc.save(input_path)
            mapping_path.write_text(json.dumps({"entries": []}), encoding="utf-8")

            with self.assertRaises(ValueError):
                restore_file(input_path=input_path, output_path=output_path, mapping_path=mapping_path)

    def test_pptx_restore_uses_confirmed_unknown_placeholder_mapping(self) -> None:
        # 失败用例：外部 AI 生成了映射表里不存在的新编号时，用户确认后的映射应能用于 PPTX 还原。
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            input_path = root / "ai_output.pptx"
            output_path = root / "restored.pptx"
            mapping_path = root / "mapping.json"

            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            box = slide.shapes.add_textbox(0, 0, 4000000, 800000)
            box.text_frame.text = "完成 PROJECT_015 与瑞为协议终稿确定。"
            prs.save(input_path)
            mapping_path.write_text(
                json.dumps(
                    {
                        "entries": [
                            {
                                "placeholder": "__PROJECT_011__",
                                "original": "瑞为协议",
                                "category": "PROJECT",
                                "enabled": True,
                            }
                        ]
                    },
                    ensure_ascii=False,
                ),
                encoding="utf-8",
            )

            restore_file(
                input_path=input_path,
                output_path=output_path,
                mapping_path=mapping_path,
                placeholder_repairs={"PROJECT_015": "__PROJECT_011__"},
            )

            restored_text = "\n".join(collect_texts_for_path(output_path))
            self.assertIn("瑞为协议", restored_text)
            self.assertNotIn("PROJECT_015", restored_text)

    def test_pptx_restore_allows_custom_unknown_placeholder_text(self) -> None:
        # 失败用例：映射表里没有对应项时，用户仍应能手动输入原词完成还原。
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            input_path = root / "ai_output.pptx"
            output_path = root / "restored.pptx"
            mapping_path = root / "mapping.json"

            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            box = slide.shapes.add_textbox(0, 0, 4000000, 800000)
            box.text_frame.text = "供应商 __SUPPLIER_004T__ 已签约。"
            prs.save(input_path)
            mapping_path.write_text(
                json.dumps(
                    {
                        "entries": [
                            {
                                "placeholder": "__SUPPLIER_001__",
                                "original": "AmpLink",
                                "category": "SUPPLIER",
                                "enabled": True,
                            }
                        ]
                    },
                    ensure_ascii=False,
                ),
                encoding="utf-8",
            )

            restore_file(
                input_path=input_path,
                output_path=output_path,
                mapping_path=mapping_path,
                placeholder_repairs={"__SUPPLIER_004T__": "Custom & Partner"},
            )

            restored_text = "\n".join(collect_texts_for_path(output_path))
            self.assertIn("Custom & Partner", restored_text)
            self.assertNotIn("__SUPPLIER_004T__", restored_text)


class LlmAssistPerformanceTests(unittest.TestCase):
    def test_llm_text_selection_filters_low_signal_texts_and_caps_chunks(self) -> None:
        # 性能回归用例：全文 XML 扫描会带来很多低价值碎片，
        # 这些碎片不应全部送给 Ollama，避免模型调用次数暴涨。
        texts = [f"普通说明文字 {idx}" for idx in range(120)]
        texts.extend(
            [
                "客户：Acme Technology Co Ltd",
                "项目：Project Phoenix",
                "合同编号：TECH-2026-001",
            ]
        )

        selected = select_texts_for_llm(texts, max_texts=20)
        chunks = chunk_texts_for_llm(selected, max_chars=30, max_items=2, max_chunks=3)

        self.assertLess(len(selected), len(texts))
        self.assertIn("客户：Acme Technology Co Ltd", selected)
        self.assertLessEqual(len(chunks), 3)

    def test_llm_candidate_contexts_skip_existing_mapping_terms(self) -> None:
        # 性能/精度用例：已有映射中的原词不需要再次交给模型判断；
        # 模型只审核规则候选和上下文，减少重复调用。
        texts = [
            "客户 Acme Technology Co Ltd 与 Project Phoenix 沟通。",
            "客户 Acme Technology Co Ltd 与 Project Phoenix 沟通。",
        ]
        rule_candidates = [
            ("COMPANY", "Acme Technology Co Ltd", "auto"),
            ("PROJECT", "Project Phoenix", "auto"),
        ]

        contexts = build_llm_candidate_contexts(texts, rule_candidates, existing_terms={"Acme Technology Co Ltd"})

        self.assertEqual(len(contexts), 1)
        self.assertNotIn("Acme Technology Co Ltd", contexts[0].split("\n", 1)[0])
        self.assertIn("Project Phoenix", contexts[0])
        self.assertEqual(stable_text_hash(contexts[0]), stable_text_hash(contexts[0]))

    def test_llm_text_selection_keeps_free_extraction_for_new_entities(self) -> None:
        # 精度回归用例：初筛必须保留原文自由抽取能力；
        # 规则没抓到的新实体仍应能进入模型输入。
        selected = select_texts_for_llm(["新出现的 Alpha Beta Legal 需要模型自由识别。"], max_texts=10)

        self.assertIn("新出现的 Alpha Beta Legal 需要模型自由识别。", selected)

    def test_existing_mapping_terms_are_counted_for_gui_logging(self) -> None:
        # GUI 日志用例：继续识别时应能看到已有映射排除了多少相关段落。
        count = count_texts_with_existing_terms(
            ["Acme 已在映射里。", "Project Phoenix 是新项目。", "Acme 再次出现。"],
            {"Acme"},
        )

        self.assertEqual(count, 2)


class UpdateCheckerTests(unittest.TestCase):
    def test_update_asset_selection_and_path_helpers(self) -> None:
        # 通过用例：同一个 Release 里有多个平台产物时，应按当前系统选择下载包。
        assets = [
            ReleaseAsset("v1.2.0-FileToolbox-macos.zip", "https://example.test/mac.zip"),
            ReleaseAsset("v1.2.0-FileToolbox.exe", "https://example.test/win.exe"),
        ]
        info = UpdateInfo(
            current_version="1.1.0",
            latest_version="1.2.0",
            release_name="v1.2.0",
            release_url="https://example.test/release",
            published_at="",
            assets=assets,
            is_update_available=True,
        )

        with patch.object(sys, "platform", "win32"):
            self.assertEqual(info.preferred_asset_name, "v1.2.0-FileToolbox.exe")
        with patch.object(sys, "platform", "darwin"):
            self.assertEqual(info.preferred_asset_name, "v1.2.0-FileToolbox-macos.zip")

        self.assertGreater(compare_versions("1.2.0", "1.1.9"), 0)
        self.assertEqual(sanitize_filename('bad:name*.exe'), "bad_name_.exe")
        with TemporaryDirectory() as temp_dir:
            existing = Path(temp_dir) / "FileToolbox.exe"
            existing.write_text("old", encoding="utf-8")
            self.assertEqual(unique_path(existing).name, "FileToolbox_2.exe")

    def test_update_download_rejects_release_without_assets(self) -> None:
        # 失败用例：Release 没有任何资产时不能假装下载成功，应提示用户打开 Release 页面人工查看。
        info = UpdateInfo(
            current_version="1.1.0",
            latest_version="1.2.0",
            release_name="v1.2.0",
            release_url="https://example.test/release",
            published_at="",
            assets=[],
            is_update_available=True,
        )

        with self.assertRaises(RuntimeError):
            download_release_asset(info)

    def test_unknown_version_string_compares_as_zero(self) -> None:
        # 失败/边界用例：无法解析数字的版本号按 0 处理，不应抛异常影响启动检测。
        self.assertEqual(compare_versions("dev", "0.0.0"), 0)
        self.assertLess(compare_versions("dev", "1.0.0"), 0)

    def test_windows_self_update_uses_new_versioned_exe_name(self) -> None:
        # 失败用例：旧逻辑把新 exe 复制到旧路径，导致内容已更新但文件名仍是旧版本号。
        current_app = Path(r"C:\Tools\v1.2.0-FileToolbox.exe")
        asset_path = Path(r"C:\Users\tester\Downloads\v1.2.1-FileToolbox.exe")

        target = windows_update_target_path(asset_path, current_app)
        script = build_windows_update_script(
            asset_path=asset_path,
            current_app=current_app,
            target_app=target,
            pid=1234,
            script_path=Path(r"C:\Temp\FileToolbox_update.bat"),
            vbs_path=Path(r"C:\Temp\FileToolbox_update.vbs"),
            log_path=Path(r"C:\Temp\FileToolbox_update.log"),
        )

        self.assertEqual(target, Path(r"C:\Tools\v1.2.1-FileToolbox.exe"))
        self.assertIn('copy /Y "%NEW_EXE%" "%TARGET_EXE%"', script)
        self.assertIn('del "%OLD_EXE%"', script)
        self.assertIn('start "" "%TARGET_EXE%"', script)
        self.assertNotIn("pause", script.lower())


class GuiLogBridgeTests(unittest.TestCase):
    def test_common_log_can_be_routed_to_gui_queue(self) -> None:
        # GUI 日志桥接用例：底层 log 仍然打印终端，同时可转发到 GUI 运行日志。
        rows: list[tuple[str, str, str]] = []
        with route_logs_to(lambda message, level, formatted: rows.append((message, level, formatted))):
            log("AI 辅助识别准备: 全文 10 段", level="INFO")

        self.assertEqual(rows[0][0], "AI 辅助识别准备: 全文 10 段")
        self.assertEqual(rows[0][1], "INFO")
        self.assertIn("[INFO] AI 辅助识别准备", rows[0][2])


if __name__ == "__main__":
    unittest.main()
