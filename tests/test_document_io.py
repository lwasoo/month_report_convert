"""Regression tests for Office document sanitize/restore I/O.

These tests exercise real DOCX/PPTX round trips because run-level text replacement and OOXML
package patching can diverge if only pure string helpers are tested.
"""

from __future__ import annotations

import json
import unittest
from tempfile import TemporaryDirectory
from pathlib import Path

from docx import Document
from pptx import Presentation

from doc_sanitizer.document_io import apply_mapping_to_file, collect_texts_for_path, restore_file
from doc_sanitizer.mapping import read_mapping


class DocumentRoundTripTests(unittest.TestCase):
    """Validate document-level behavior at the public file I/O boundary."""

    def test_docx_sanitize_and_restore_with_confirmed_placeholder_repairs(self) -> None:
        # 通过用例：覆盖最核心的 docx 链路：
        # 原文 -> 脱敏文档/映射 JSON -> 外部 AI 改坏占位符 -> 按确认结果还原。
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            source_path = root / "source.docx"
            sanitized_path = root / "sanitized.docx"
            edited_path = root / "edited.docx"
            restored_path = root / "restored.docx"
            mapping_path = root / "mapping.json"

            doc = Document()
            doc.add_paragraph("Acme Technology Co Ltd signed Project Phoenix.")
            doc.save(source_path)

            payload = {
                "version": 2,
                "entries": [
                    {
                        "placeholder": "__COMPANY_001__",
                        "original": "Acme Technology Co Ltd",
                        "category": "COMPANY",
                        "enabled": True,
                        "source": "test",
                    },
                    {
                        "placeholder": "__PROJECT_001__",
                        "original": "Project Phoenix",
                        "category": "PROJECT",
                        "enabled": True,
                        "source": "test",
                    },
                ],
            }

            apply_mapping_to_file(source_path, sanitized_path, payload, mapping_path)
            sanitized_text = "\n".join(p.text for p in Document(sanitized_path).paragraphs)
            self.assertIn("__COMPANY_001__", sanitized_text)
            self.assertIn("__PROJECT_001__", sanitized_text)
            self.assertEqual(read_mapping(mapping_path)["replacements"]["__COMPANY_001__"], "Acme Technology Co Ltd")

            edited = Document()
            edited.add_paragraph("COMPANY_001 signed __PROJECT-001__.")
            edited.save(edited_path)

            restore_file(
                input_path=edited_path,
                output_path=restored_path,
                mapping_path=mapping_path,
                placeholder_repairs={
                    "COMPANY_001": "__COMPANY_001__",
                    "__PROJECT-001__": "__PROJECT_001__",
                },
            )
            restored_text = "\n".join(p.text for p in Document(restored_path).paragraphs)
            self.assertIn("Acme Technology Co Ltd", restored_text)
            self.assertIn("Project Phoenix", restored_text)

    def test_restore_rejects_empty_mapping_file(self) -> None:
        # 失败用例：映射 JSON 没有有效条目时，不能生成一个看似成功但未还原的文件。
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            input_path = root / "input.docx"
            output_path = root / "output.docx"
            mapping_path = root / "mapping.json"

            doc = Document()
            doc.add_paragraph("__COMPANY_001__")
            doc.save(input_path)
            mapping_path.write_text(json.dumps({"entries": []}), encoding="utf-8")

            with self.assertRaises(ValueError):
                restore_file(input_path=input_path, output_path=output_path, mapping_path=mapping_path)

    def test_pptx_restore_uses_confirmed_unknown_placeholder_mapping(self) -> None:
        # 失败用例：外部 AI 生成了映射表里不存在的新编号时，用户确认后的映射应能用于 PPTX 还原。
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            input_path = root / "ai_output.pptx"
            output_path = root / "restored.pptx"
            mapping_path = root / "mapping.json"

            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            box = slide.shapes.add_textbox(0, 0, 4000000, 800000)
            box.text_frame.text = "完成 PROJECT_015 与瑞为协议终稿确定。"
            prs.save(input_path)
            mapping_path.write_text(
                json.dumps(
                    {
                        "entries": [
                            {
                                "placeholder": "__PROJECT_011__",
                                "original": "瑞为协议",
                                "category": "PROJECT",
                                "enabled": True,
                            }
                        ]
                    },
                    ensure_ascii=False,
                ),
                encoding="utf-8",
            )

            restore_file(
                input_path=input_path,
                output_path=output_path,
                mapping_path=mapping_path,
                placeholder_repairs={"PROJECT_015": "__PROJECT_011__"},
            )

            restored_text = "\n".join(collect_texts_for_path(output_path))
            self.assertIn("瑞为协议", restored_text)
            self.assertNotIn("PROJECT_015", restored_text)

    def test_pptx_restore_allows_custom_unknown_placeholder_text(self) -> None:
        # 失败用例：映射表里没有对应项时，用户仍应能手动输入原词完成还原。
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            input_path = root / "ai_output.pptx"
            output_path = root / "restored.pptx"
            mapping_path = root / "mapping.json"

            prs = Presentation()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            box = slide.shapes.add_textbox(0, 0, 4000000, 800000)
            box.text_frame.text = "供应商 __SUPPLIER_004T__ 已签约。"
            prs.save(input_path)
            mapping_path.write_text(
                json.dumps(
                    {
                        "entries": [
                            {
                                "placeholder": "__SUPPLIER_001__",
                                "original": "AmpLink",
                                "category": "SUPPLIER",
                                "enabled": True,
                            }
                        ]
                    },
                    ensure_ascii=False,
                ),
                encoding="utf-8",
            )

            restore_file(
                input_path=input_path,
                output_path=output_path,
                mapping_path=mapping_path,
                placeholder_repairs={"__SUPPLIER_004T__": "Custom & Partner"},
            )

            restored_text = "\n".join(collect_texts_for_path(output_path))
            self.assertIn("Custom & Partner", restored_text)
            self.assertNotIn("__SUPPLIER_004T__", restored_text)


if __name__ == "__main__":
    unittest.main()
