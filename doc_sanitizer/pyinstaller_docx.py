"""Compatibility fixes for Office libraries inside PyInstaller app bundles."""

from __future__ import annotations

import sys
from pathlib import Path


def _resource_roots() -> list[Path]:
    roots: list[Path] = []
    meipass = getattr(sys, "_MEIPASS", None)
    if meipass:
        root = Path(meipass)
        roots.extend(
            [
                root,
                root.parent / "Resources",
                root.parent / "Frameworks",
            ]
        )

    executable = Path(sys.executable).resolve()
    if executable.name:
        contents = executable.parent.parent
        roots.extend([contents / "Resources", contents / "Frameworks"])

    unique_roots: list[Path] = []
    seen: set[Path] = set()
    for root in roots:
        resolved = root.resolve()
        if resolved not in seen:
            unique_roots.append(resolved)
            seen.add(resolved)
    return unique_roots


def _template_candidates(package: str, filename: str) -> list[Path]:
    return [root / package / "templates" / filename for root in _resource_roots()]


def _read_template(package: str, filename: str) -> bytes:
    candidates = _template_candidates(package, filename)
    for candidate in candidates:
        if candidate.exists():
            return candidate.read_bytes()
    checked = ", ".join(str(candidate) for candidate in candidates) or "<no PyInstaller roots>"
    raise FileNotFoundError(f"Unable to find {package} template {filename}; checked: {checked}")


def _template_path(package: str, filename: str) -> str:
    candidates = _template_candidates(package, filename)
    for candidate in candidates:
        if candidate.exists():
            return str(candidate)
    checked = ", ".join(str(candidate) for candidate in candidates) or "<no PyInstaller roots>"
    raise FileNotFoundError(f"Unable to find {package} template {filename}; checked: {checked}")


def patch_python_docx_templates() -> None:
    """Make python-docx find bundled templates in macOS PyInstaller app layout."""
    if not getattr(sys, "frozen", False):
        return

    from docx.parts.hdrftr import FooterPart, HeaderPart

    original_header = HeaderPart._default_header_xml
    original_footer = FooterPart._default_footer_xml

    def default_header_xml(cls) -> bytes:
        try:
            return original_header()
        except FileNotFoundError:
            return _read_template("docx", "default-header.xml")

    def default_footer_xml(cls) -> bytes:
        try:
            return original_footer()
        except FileNotFoundError:
            return _read_template("docx", "default-footer.xml")

    HeaderPart._default_header_xml = classmethod(default_header_xml)
    FooterPart._default_footer_xml = classmethod(default_footer_xml)


def patch_python_pptx_templates() -> None:
    """Make python-pptx find bundled templates in macOS PyInstaller app layout."""
    if not getattr(sys, "frozen", False):
        return

    import pptx.api
    import pptx.oxml
    import pptx.shapes.shapetree
    from pptx.util import lazyproperty

    original_default_pptx_path = pptx.api._default_pptx_path
    original_parse_from_template = pptx.oxml.parse_from_template
    original_icon_image_file = pptx.shapes.shapetree._OleObjectElementCreator._icon_image_file

    def default_pptx_path() -> str:
        path = Path(original_default_pptx_path())
        if path.exists():
            return str(path)
        return _template_path("pptx", "default.pptx")

    def parse_from_template(template_file_name: str):
        try:
            return original_parse_from_template(template_file_name)
        except FileNotFoundError:
            xml = _read_template("pptx", f"{template_file_name}.xml")
            return pptx.oxml.parse_xml(xml)

    def icon_image_file(self):
        try:
            path = original_icon_image_file.__get__(self, type(self))
            if not isinstance(path, str) or Path(path).exists():
                return path
        except FileNotFoundError:
            pass

        icon_filename = (
            self._prog_id_arg.icon_filename
            if isinstance(self._prog_id_arg, pptx.shapes.shapetree.PROG_ID)
            else "generic-icon.emf"
        )
        return _template_path("pptx", icon_filename)

    pptx.api._default_pptx_path = default_pptx_path
    pptx.oxml.parse_from_template = parse_from_template
    pptx.shapes.shapetree._OleObjectElementCreator._icon_image_file = lazyproperty(icon_image_file)


def patch_office_library_templates() -> None:
    """Apply all Office-library PyInstaller template path compatibility fixes."""
    patch_python_docx_templates()
    patch_python_pptx_templates()
