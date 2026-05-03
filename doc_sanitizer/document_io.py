from __future__ import annotations

from pathlib import Path
from tempfile import TemporaryDirectory
from typing import Any
from xml.sax.saxutils import escape
from xml.etree import ElementTree as ET
from zipfile import ZIP_DEFLATED, ZipFile

from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from office_conversion import convert_from_ooxml, convert_to_ooxml
from report_converter.common import log, normalize_text
from .fuzzy_mapping import repair_placeholder_text
from .mapping import ReplacementItem, mapping_entries, read_mapping, write_mapping_data

SUPPORTED_FILE_SUFFIXES = {".doc", ".docx", ".ppt", ".pptx"}


def ensure_supported_path(input_path: Path) -> None:
    if input_path.suffix.lower() not in SUPPORTED_FILE_SUFFIXES:
        raise ValueError("当前仅支持 .doc/.docx 和 .ppt/.pptx 文件。")


def default_sanitized_path(input_path: Path) -> Path:
    ensure_supported_path(input_path)
    return input_path.with_name(f"{input_path.stem}_脱敏{input_path.suffix.lower()}")


def collect_texts_for_path(input_path: Path) -> list[str]:
    ensure_supported_path(input_path)
    if input_path.suffix.lower() in {".doc", ".ppt"}:
        with TemporaryDirectory() as td:
            converted = convert_to_ooxml(input_path, Path(td))
            return collect_texts_for_path(converted)
    if input_path.suffix.lower() == ".docx":
        return dedupe_texts([*collect_doc_texts(Document(str(input_path))), *collect_docx_package_texts(input_path)])
    return dedupe_texts([*collect_ppt_texts(Presentation(str(input_path))), *collect_pptx_package_texts(input_path)])


def collect_doc_texts(doc: Document) -> list[str]:
    texts: list[str] = []
    for paragraph in iter_doc_paragraphs(doc):
        text = normalize_text(paragraph.text)
        if text:
            texts.append(text)
    return texts


def dedupe_texts(texts: list[str]) -> list[str]:
    out: list[str] = []
    seen: set[str] = set()
    for text in texts:
        normalized = normalize_text(text)
        if not normalized or normalized in seen:
            continue
        seen.add(normalized)
        out.append(normalized)
    return out


def collect_docx_package_texts(input_path: Path) -> list[str]:
    texts: list[str] = []
    with ZipFile(input_path, "r") as zin:
        for name in zin.namelist():
            if not is_docx_text_xml_part(name):
                continue
            try:
                root = ET.fromstring(zin.read(name))
            except ET.ParseError:
                continue
            for elem in root.iter():
                if elem.text:
                    text = normalize_text(elem.text)
                    if text:
                        texts.append(text)
    return texts


def collect_pptx_package_texts(input_path: Path) -> list[str]:
    texts: list[str] = []
    with ZipFile(input_path, "r") as zin:
        for name in zin.namelist():
            if not is_pptx_text_xml_part(name):
                continue
            try:
                root = ET.fromstring(zin.read(name))
            except ET.ParseError:
                continue
            for elem in root.iter():
                if elem.text:
                    text = normalize_text(elem.text)
                    if text:
                        texts.append(text)
    return texts


def is_docx_text_xml_part(name: str) -> bool:
    if not name.startswith("word/") or not name.endswith(".xml"):
        return False
    if name.startswith("word/_rels/"):
        return False
    return True


def is_pptx_text_xml_part(name: str) -> bool:
    if not name.startswith("ppt/") or not name.endswith(".xml"):
        return False
    if "/_rels/" in name or name.startswith("ppt/_rels/"):
        return False
    return True


def collect_ppt_texts(prs: Presentation) -> list[str]:
    texts: list[str] = []
    for paragraph in iter_ppt_paragraphs(prs):
        text = normalize_text(paragraph.text)
        if text:
            texts.append(text)
    return texts


def iter_doc_paragraphs(doc: Document):
    for paragraph in doc.paragraphs:
        yield paragraph
    for table in doc.tables:
        yield from iter_doc_table_paragraphs(table)
    for section in doc.sections:
        for paragraph in section.header.paragraphs:
            yield paragraph
        for table in section.header.tables:
            yield from iter_doc_table_paragraphs(table)
        for paragraph in section.footer.paragraphs:
            yield paragraph
        for table in section.footer.tables:
            yield from iter_doc_table_paragraphs(table)


def iter_doc_table_paragraphs(table):
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.paragraphs:
                yield paragraph
            for inner in cell.tables:
                yield from iter_doc_table_paragraphs(inner)


def iter_ppt_paragraphs(prs: Presentation):
    for slide in prs.slides:
        for shape in slide.shapes:
            yield from iter_ppt_shape_paragraphs(shape)


def iter_ppt_shape_paragraphs(shape):
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        for inner in shape.shapes:
            yield from iter_ppt_shape_paragraphs(inner)
        return
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    yield paragraph
        return
    if getattr(shape, "has_text_frame", False):
        for paragraph in shape.text_frame.paragraphs:
            yield paragraph


def apply_replacements_to_doc(
    doc: Document,
    items: list[ReplacementItem],
    reverse: bool = False,
    placeholder_repairs: dict[str, str] | None = None,
) -> None:
    ordered = sorted(items, key=lambda item: len(item.placeholder if reverse else item.original), reverse=True)
    for paragraph in iter_doc_paragraphs(doc):
        replace_in_doc_paragraph(paragraph, ordered, reverse=reverse, placeholder_repairs=placeholder_repairs)


def apply_replacements_to_ppt(
    prs: Presentation,
    items: list[ReplacementItem],
    reverse: bool = False,
    placeholder_repairs: dict[str, str] | None = None,
) -> None:
    ordered = sorted(items, key=lambda item: len(item.placeholder if reverse else item.original), reverse=True)
    for paragraph in iter_ppt_paragraphs(prs):
        replace_in_ppt_paragraph(paragraph, ordered, reverse=reverse, placeholder_repairs=placeholder_repairs)


def replace_in_doc_paragraph(
    paragraph,
    items: list[ReplacementItem],
    reverse: bool = False,
    placeholder_repairs: dict[str, str] | None = None,
) -> None:
    source = paragraph.text or ""
    if not source:
        return
    if replace_in_runs(paragraph.runs, items, reverse=reverse, placeholder_repairs=placeholder_repairs):
        return
    updated = replace_text(source, items, reverse=reverse, placeholder_repairs=placeholder_repairs)
    if updated == source:
        return
    if len(paragraph.runs) == 1:
        paragraph.runs[0].text = updated
        return
    if paragraph.runs:
        paragraph.runs[0].text = updated
        for run in paragraph.runs[1:]:
            run.text = ""
        return
    paragraph.add_run(updated)


def replace_in_ppt_paragraph(
    paragraph,
    items: list[ReplacementItem],
    reverse: bool = False,
    placeholder_repairs: dict[str, str] | None = None,
) -> None:
    source = paragraph.text or ""
    if not source:
        return
    runs = list(paragraph.runs)
    if replace_in_runs(runs, items, reverse=reverse, placeholder_repairs=placeholder_repairs):
        return
    updated = replace_text(source, items, reverse=reverse, placeholder_repairs=placeholder_repairs)
    if updated == source:
        return
    if len(runs) == 1:
        runs[0].text = updated
        return
    if runs:
        runs[0].text = updated
        for run in runs[1:]:
            run.text = ""
        return
    paragraph.text = updated


def replace_text(
    text: str,
    items: list[ReplacementItem],
    reverse: bool = False,
    placeholder_repairs: dict[str, str] | None = None,
) -> str:
    updated = repair_placeholder_text(text, items, confirmed_repairs=placeholder_repairs) if reverse else text
    for item in items:
        old = item.placeholder if reverse else item.original
        new = item.original if reverse else item.placeholder
        if old in updated:
            updated = updated.replace(old, new)
    return updated


def replace_in_runs(
    runs,
    items: list[ReplacementItem],
    reverse: bool = False,
    placeholder_repairs: dict[str, str] | None = None,
) -> bool:
    run_list = list(runs)
    if not run_list:
        return False
    texts = [run.text or "" for run in run_list]
    source = "".join(texts)
    if not source:
        return False
    if reverse and placeholder_repairs:
        return False

    char_runs: list[int] = []
    for run_idx, text in enumerate(texts):
        char_runs.extend([run_idx] * len(text))

    chunks: list[list[str]] = [[] for _ in run_list]
    changed = False
    pos = 0
    while pos < len(source):
        matched = False
        for item in items:
            old = item.placeholder if reverse else item.original
            new = item.original if reverse else item.placeholder
            if old and source.startswith(old, pos):
                chunks[char_runs[pos]].append(new)
                pos += len(old)
                changed = True
                matched = True
                break
        if matched:
            continue
        chunks[char_runs[pos]].append(source[pos])
        pos += 1

    if not changed:
        return False
    for run, parts in zip(run_list, chunks):
        run.text = "".join(parts)
    return True


def apply_replacements_to_docx_package(
    path: Path,
    items: list[ReplacementItem],
    reverse: bool = False,
    placeholder_repairs: dict[str, str] | None = None,
) -> None:
    if not items:
        return
    temp_path = path.with_name(f"{path.stem}.tmp{path.suffix}")
    changed_any = False
    with ZipFile(path, "r") as zin, ZipFile(temp_path, "w", ZIP_DEFLATED) as zout:
        for info in zin.infolist():
            data = zin.read(info.filename)
            if is_docx_text_xml_part(info.filename):
                try:
                    text = data.decode("utf-8")
                except UnicodeDecodeError:
                    text = ""
                if text:
                    updated = replace_xml_text(text, items, reverse=reverse, placeholder_repairs=placeholder_repairs)
                    if updated != text:
                        data = updated.encode("utf-8")
                        changed_any = True
            zout.writestr(info, data)
    if changed_any:
        temp_path.replace(path)
    else:
        temp_path.unlink(missing_ok=True)


def apply_replacements_to_pptx_package(
    path: Path,
    items: list[ReplacementItem],
    reverse: bool = False,
    placeholder_repairs: dict[str, str] | None = None,
) -> None:
    if not items:
        return
    temp_path = path.with_name(f"{path.stem}.tmp{path.suffix}")
    changed_any = False
    with ZipFile(path, "r") as zin, ZipFile(temp_path, "w", ZIP_DEFLATED) as zout:
        for info in zin.infolist():
            data = zin.read(info.filename)
            if is_pptx_text_xml_part(info.filename):
                try:
                    text = data.decode("utf-8")
                except UnicodeDecodeError:
                    text = ""
                if text:
                    updated = replace_xml_text(text, items, reverse=reverse, placeholder_repairs=placeholder_repairs)
                    if updated != text:
                        data = updated.encode("utf-8")
                        changed_any = True
            zout.writestr(info, data)
    if changed_any:
        temp_path.replace(path)
    else:
        temp_path.unlink(missing_ok=True)


def replace_xml_text(
    xml_text: str,
    items: list[ReplacementItem],
    reverse: bool = False,
    placeholder_repairs: dict[str, str] | None = None,
) -> str:
    xml_repairs = {token: escape(value) for token, value in (placeholder_repairs or {}).items()}
    updated = repair_placeholder_text(xml_text, items, confirmed_repairs=xml_repairs) if reverse else xml_text
    for item in items:
        old = item.placeholder if reverse else item.original
        new = item.original if reverse else item.placeholder
        if old:
            updated = updated.replace(escape(old), escape(new))
    return updated


def apply_mapping_to_file(
    input_path: Path,
    output_path: Path,
    payload: dict[str, Any],
    mapping_path: Path | None = None,
) -> None:
    ensure_supported_path(input_path)
    if input_path.suffix.lower() in {".doc", ".ppt"} or output_path.suffix.lower() in {".doc", ".ppt"}:
        with TemporaryDirectory() as td:
            work_dir = Path(td)
            converted_input = convert_to_ooxml(input_path, work_dir)
            temp_output = work_dir / f"{output_path.stem}{converted_input.suffix.lower()}"
            apply_mapping_to_file(converted_input, temp_output, payload, None)
            convert_from_ooxml(temp_output, output_path)
        if mapping_path is not None:
            mapping_path.parent.mkdir(parents=True, exist_ok=True)
            payload["sanitized_file"] = str(output_path)
            write_mapping_data(mapping_path, payload)
        return

    if input_path.suffix.lower() == ".docx":
        items = mapping_entries(payload)
        doc = Document(str(input_path))
        apply_replacements_to_doc(doc, items)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(output_path))
        apply_replacements_to_docx_package(output_path, items)
    else:
        prs = Presentation(str(input_path))
        apply_replacements_to_ppt(prs, mapping_entries(payload))
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        apply_replacements_to_pptx_package(output_path, mapping_entries(payload))
    if mapping_path is not None:
        mapping_path.parent.mkdir(parents=True, exist_ok=True)
        payload["sanitized_file"] = str(output_path)
        write_mapping_data(mapping_path, payload)


def restore_file(
    input_path: Path,
    output_path: Path,
    mapping_path: Path,
    placeholder_repairs: dict[str, str] | None = None,
) -> None:
    ensure_supported_path(input_path)
    payload = read_mapping(mapping_path)
    items = mapping_entries(payload, only_enabled=False)
    if not items:
        raise ValueError("映射文件中未找到有效 entries。")
    if input_path.suffix.lower() in {".doc", ".ppt"} or output_path.suffix.lower() in {".doc", ".ppt"}:
        with TemporaryDirectory() as td:
            work_dir = Path(td)
            converted_input = convert_to_ooxml(input_path, work_dir)
            temp_output = work_dir / f"{output_path.stem}{converted_input.suffix.lower()}"
            restore_file(converted_input, temp_output, mapping_path, placeholder_repairs=placeholder_repairs)
            convert_from_ooxml(temp_output, output_path)
        log(f"还原完成: {output_path}")
        return

    if input_path.suffix.lower() == ".docx":
        doc = Document(str(input_path))
        apply_replacements_to_doc(doc, items, reverse=True, placeholder_repairs=placeholder_repairs)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(output_path))
        apply_replacements_to_docx_package(output_path, items, reverse=True, placeholder_repairs=placeholder_repairs)
    else:
        prs = Presentation(str(input_path))
        apply_replacements_to_ppt(prs, items, reverse=True, placeholder_repairs=placeholder_repairs)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        apply_replacements_to_pptx_package(output_path, items, reverse=True, placeholder_repairs=placeholder_repairs)
    log(f"还原完成: {output_path}")


def apply_mapping_to_docx(
    input_path: Path,
    output_path: Path,
    payload: dict[str, Any],
    mapping_path: Path | None = None,
) -> None:
    apply_mapping_to_file(input_path, output_path, payload, mapping_path)


def restore_docx(input_path: Path, output_path: Path, mapping_path: Path) -> None:
    restore_file(input_path, output_path, mapping_path)
