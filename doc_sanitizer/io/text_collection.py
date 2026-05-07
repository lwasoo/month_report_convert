"""Text collection helpers for Office documents.

The scanner needs the same text coverage that replacement later touches: visible paragraph
text plus XML text inside the OOXML package that high-level libraries can miss.
"""

from __future__ import annotations

from pathlib import Path
from tempfile import TemporaryDirectory
from xml.etree import ElementTree as ET
from zipfile import ZipFile

from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from office_conversion import convert_to_ooxml
from report_converter.common import normalize_text

from .file_types import ensure_supported_path, is_legacy_office_path
from .ooxml_parts import is_docx_text_xml_part, is_pptx_text_xml_part


def collect_texts_for_path(input_path: Path) -> list[str]:
    """Collect visible and package-level text for DOC/DOCX/PPT/PPTX inputs."""
    ensure_supported_path(input_path)
    if is_legacy_office_path(input_path):
        with TemporaryDirectory() as td:
            converted = convert_to_ooxml(input_path, Path(td))
            return collect_texts_for_path(converted)
    if input_path.suffix.lower() == ".docx":
        return dedupe_texts([*collect_doc_texts(Document(str(input_path))), *collect_docx_package_texts(input_path)])
    return dedupe_texts([*collect_ppt_texts(Presentation(str(input_path))), *collect_pptx_package_texts(input_path)])


def collect_doc_texts(doc: Document) -> list[str]:
    texts: list[str] = []
    for paragraph in iter_doc_paragraphs(doc):
        text = normalize_text(paragraph.text)
        if text:
            texts.append(text)
    return texts


def collect_ppt_texts(prs: Presentation) -> list[str]:
    texts: list[str] = []
    for paragraph in iter_ppt_paragraphs(prs):
        text = normalize_text(paragraph.text)
        if text:
            texts.append(text)
    return texts


def collect_docx_package_texts(input_path: Path) -> list[str]:
    """Read text nodes from DOCX XML parts that python-docx may not expose."""
    return collect_package_texts(input_path, is_docx_text_xml_part)


def collect_pptx_package_texts(input_path: Path) -> list[str]:
    """Read text nodes from PPTX XML parts that python-pptx may not expose."""
    return collect_package_texts(input_path, is_pptx_text_xml_part)


def collect_package_texts(input_path: Path, part_filter) -> list[str]:
    texts: list[str] = []
    with ZipFile(input_path, "r") as zin:
        for name in zin.namelist():
            if not part_filter(name):
                continue
            try:
                root = ET.fromstring(zin.read(name))
            except ET.ParseError:
                continue
            for elem in root.iter():
                if elem.text:
                    text = normalize_text(elem.text)
                    if text:
                        texts.append(text)
    return texts


def dedupe_texts(texts: list[str]) -> list[str]:
    out: list[str] = []
    seen: set[str] = set()
    for text in texts:
        normalized = normalize_text(text)
        if not normalized or normalized in seen:
            continue
        seen.add(normalized)
        out.append(normalized)
    return out


def iter_doc_paragraphs(doc: Document):
    for paragraph in doc.paragraphs:
        yield paragraph
    for table in doc.tables:
        yield from iter_doc_table_paragraphs(table)
    for section in doc.sections:
        for paragraph in section.header.paragraphs:
            yield paragraph
        for table in section.header.tables:
            yield from iter_doc_table_paragraphs(table)
        for paragraph in section.footer.paragraphs:
            yield paragraph
        for table in section.footer.tables:
            yield from iter_doc_table_paragraphs(table)


def iter_doc_table_paragraphs(table):
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.paragraphs:
                yield paragraph
            for inner in cell.tables:
                yield from iter_doc_table_paragraphs(inner)


def iter_ppt_paragraphs(prs: Presentation):
    for slide in prs.slides:
        for shape in slide.shapes:
            yield from iter_ppt_shape_paragraphs(shape)


def iter_ppt_shape_paragraphs(shape):
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        for inner in shape.shapes:
            yield from iter_ppt_shape_paragraphs(inner)
        return
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    yield paragraph
        return
    if getattr(shape, "has_text_frame", False):
        for paragraph in shape.text_frame.paragraphs:
            yield paragraph
