"""High-level sanitize and restore file operations."""

from __future__ import annotations

from pathlib import Path
from tempfile import TemporaryDirectory

from docx import Document
from pptx import Presentation

from office_conversion import convert_from_ooxml, convert_to_ooxml
from report_converter.common import log

from .file_types import ensure_supported_path, is_legacy_office_path
from ..mapping import MappingLike, mapping_entries, read_mapping, write_mapping_data
from .ooxml_package import apply_replacements_to_docx_package, apply_replacements_to_pptx_package
from .replacement_engine import ReplacementDirection, apply_replacements_to_doc, apply_replacements_to_ppt


def apply_mapping_to_file(
    input_path: Path,
    output_path: Path,
    payload: MappingLike,
    mapping_path: Path | None = None,
) -> None:
    """Apply a mapping to Office files, converting legacy .doc/.ppt through OOXML first."""
    ensure_supported_path(input_path)
    if is_legacy_office_path(input_path) or is_legacy_office_path(output_path):
        with TemporaryDirectory() as td:
            work_dir = Path(td)
            converted_input = convert_to_ooxml(input_path, work_dir)
            temp_output = work_dir / f"{output_path.stem}{converted_input.suffix.lower()}"
            apply_mapping_to_file(converted_input, temp_output, payload, None)
            convert_from_ooxml(temp_output, output_path)
        if mapping_path is not None:
            mapping_path.parent.mkdir(parents=True, exist_ok=True)
            payload["sanitized_file"] = str(output_path)
            write_mapping_data(mapping_path, payload)
        return

    items = mapping_entries(payload)
    if input_path.suffix.lower() == ".docx":
        doc = Document(str(input_path))
        apply_replacements_to_doc(doc, items, direction=ReplacementDirection.SANITIZE)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(output_path))
        apply_replacements_to_docx_package(output_path, items, direction=ReplacementDirection.SANITIZE)
    else:
        prs = Presentation(str(input_path))
        apply_replacements_to_ppt(prs, items, direction=ReplacementDirection.SANITIZE)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        apply_replacements_to_pptx_package(output_path, items, direction=ReplacementDirection.SANITIZE)
    if mapping_path is not None:
        mapping_path.parent.mkdir(parents=True, exist_ok=True)
        payload["sanitized_file"] = str(output_path)
        write_mapping_data(mapping_path, payload)


def restore_file(
    input_path: Path,
    output_path: Path,
    mapping_path: Path,
    placeholder_repairs: dict[str, str] | None = None,
) -> None:
    """Restore placeholders back to originals, optionally applying confirmed repairs first."""
    ensure_supported_path(input_path)
    payload = read_mapping(mapping_path)
    items = mapping_entries(payload, only_enabled=False)
    if not items:
        raise ValueError("映射文件中未找到有效 entries。")
    if is_legacy_office_path(input_path) or is_legacy_office_path(output_path):
        with TemporaryDirectory() as td:
            work_dir = Path(td)
            converted_input = convert_to_ooxml(input_path, work_dir)
            temp_output = work_dir / f"{output_path.stem}{converted_input.suffix.lower()}"
            restore_file(converted_input, temp_output, mapping_path, placeholder_repairs=placeholder_repairs)
            convert_from_ooxml(temp_output, output_path)
        log(f"还原完成: {output_path}")
        return

    if input_path.suffix.lower() == ".docx":
        doc = Document(str(input_path))
        apply_replacements_to_doc(
            doc,
            items,
            direction=ReplacementDirection.RESTORE,
            placeholder_repairs=placeholder_repairs,
        )
        output_path.parent.mkdir(parents=True, exist_ok=True)
        doc.save(str(output_path))
        apply_replacements_to_docx_package(
            output_path,
            items,
            direction=ReplacementDirection.RESTORE,
            placeholder_repairs=placeholder_repairs,
        )
    else:
        prs = Presentation(str(input_path))
        apply_replacements_to_ppt(
            prs,
            items,
            direction=ReplacementDirection.RESTORE,
            placeholder_repairs=placeholder_repairs,
        )
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        apply_replacements_to_pptx_package(
            output_path,
            items,
            direction=ReplacementDirection.RESTORE,
            placeholder_repairs=placeholder_repairs,
        )
    log(f"还原完成: {output_path}")


def apply_mapping_to_docx(
    input_path: Path,
    output_path: Path,
    payload: MappingLike,
    mapping_path: Path | None = None,
) -> None:
    apply_mapping_to_file(input_path, output_path, payload, mapping_path)


def restore_docx(input_path: Path, output_path: Path, mapping_path: Path) -> None:
    restore_file(input_path, output_path, mapping_path)
